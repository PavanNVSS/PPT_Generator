import os
from pptx import Presentation
from win32com.client import Dispatch

def save_slide_layouts_as_images(designs_folder, output_folder):
    """
    Save each slide layout from PowerPoint templates as images using PowerPoint COM API.

    Args:
        designs_folder (str): Folder containing PowerPoint design templates.
        output_folder (str): Folder to save the layout images.
    """
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    design_files = [f for f in os.listdir(designs_folder) if f.endswith('.pptx')]

    # Start PowerPoint application
    powerpoint = Dispatch("PowerPoint.Application")
    powerpoint.Visible = 1

    for design_file in design_files:
        design_path = os.path.abspath(os.path.join(designs_folder, design_file))
        prs = Presentation(design_path)

        for index, layout in enumerate(prs.slide_layouts):
            # Create a temporary PowerPoint file with the current layout
            temp_ppt = Presentation()
            slide = temp_ppt.slides.add_slide(layout)

            # Add sample content for layout preview
            if slide.shapes.title:
                slide.shapes.title.text = f"Layout {index + 1}"
            for shape in slide.placeholders:
                if shape.placeholder_format.idx != 0:  # Exclude title
                    shape.text = f"Sample Content for Layout {index + 1}"

            temp_pptx_path = os.path.abspath(os.path.join(output_folder, f"temp_{index + 1}.pptx"))
            temp_ppt.save(temp_pptx_path)

            # Verify the temporary file exists
            if not os.path.exists(temp_pptx_path):
                print(f"Temporary file not created: {temp_pptx_path}")
                continue

            # Open the temporary PowerPoint in PowerPoint application
            try:
                presentation = powerpoint.Presentations.Open(temp_pptx_path)
                output_image_path = os.path.abspath(os.path.join(output_folder, f"{os.path.splitext(design_file)[0]}_layout_{index + 1}.jpg"))

                # Export the first slide as an image
                presentation.Slides[1].Export(output_image_path, "JPG")
                presentation.Close()

                print(f"Saved layout {index + 1} for {design_file} as an image.")
            except Exception as e:
                print(f"Error exporting layout {index + 1} for {design_file}: {e}")
                continue
            finally:
                # Remove the temporary PowerPoint file
                if os.path.exists(temp_pptx_path):
                    os.remove(temp_pptx_path)

    # Quit PowerPoint application
    powerpoint.Quit()

# Example usage
designs_folder = "Designs"
output_folder = "LayoutPreviews"
save_slide_layouts_as_images(designs_folder, output_folder)
