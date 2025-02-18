# powerpoint.py
import os
import random
import re
from pptx import Presentation


def create_ppt_text(prompt, slides, add_info):
    """
    Generate the text content for the PowerPoint.

    Args:
        prompt (str): The topic of the presentation.
        slides (int): The number of slides.
        add_info (str): Additional information or instructions.

    Returns:
        str: The formatted text content for the PowerPoint.
    """

    print("this is before processing of text in ppt : ",prompt)

    content = []
    content.append(f"Title: {prompt}")
    # content.append(f"Header: Introduction")
    # content.append(f"Content: This is a presentation on {prompt}. {add_info}\n")

    # for i in range(1, int(slides) + 1):
    #     content.append(f"Slide: {i}")
    #     content.append(f"Header: Slide {i}")
    #     content.append(f"Content: This is slide {i} content. Insert detailed information here.")

    print("this is after processing of text in ppt : ",prompt)

    return "\n".join(content)


def clean_text_file(text_file):
    """
    Cleans the input text file by removing recurring content such as duplicate
    'Header: Introduction' and 'Content: This is a presentation on...' entries.

    Args:
        text_file (str): Path to the text file.

    Returns:
        str: Cleaned content as a single string.
    """
    intro_content_pattern = "This is a presentation on"
    header_count = 0
    content_count = 0

    cleaned_lines = []
    with open(text_file, 'r', encoding='utf-8') as f:
        for line in f:
            # Remove duplicate "Header: Introduction"
            if line.startswith("Header: Introduction"):
                header_count += 1
                if header_count > 1:
                    continue

            # Remove duplicate introduction content
            if intro_content_pattern in line:
                content_count += 1
                if content_count > 1:
                    continue

            cleaned_lines.append(line.strip())

    return "\n".join(cleaned_lines)




def create_ppt(text_file, design_number, ppt_name, slide_count):
    prs = Presentation(f"Designs/Design-{design_number}.pptx")
    slide_count_created = 0
    header = ""
    content = ""
    last_slide_layout_index = -1
    firsttime = True
    clean_text = clean_text_file(text_file)
    print("this is clean text",clean_text)

    with open(text_file, 'r', encoding='utf-8') as f:
        for line_num, line in enumerate(f):
            if slide_count_created >= slide_count:
                break


            if line.startswith('Title:'):

                header = line.replace('Title:', '').strip()
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                title.text = header
                body_shape = slide.shapes.placeholders[1]
                continue
            elif line.startswith('Slide:'):
                if slide_count_created > 0:

                    slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
                    title = slide.shapes.title
                    title.text = header
                    body_shape = slide.shapes.placeholders[slide_placeholder_index]
                    tf = body_shape.text_frame
                    tf.text = content


                content = ""
                slide_count_created += 1
                slide_layout_index = last_slide_layout_index
                layout_indices = [1, 7, 8]
                while slide_layout_index == last_slide_layout_index:
                    if firsttime == True:
                        slide_layout_index = 1
                        slide_placeholder_index = 1
                        firsttime = False
                        break
                    slide_layout_index = random.choice(layout_indices)
                    if slide_layout_index == 8:
                        slide_placeholder_index = 2
                    else:
                        slide_placeholder_index = 1
                last_slide_layout_index = slide_layout_index
                continue
            elif line.startswith('Header:'):


                header = line.replace('Header:', '').strip()
                continue
            elif line.startswith('Content:'):

                content = line.replace('Content:', '').strip()
                next_line = f.readline().strip()
                while next_line and not next_line.startswith('#'):
                    content += '\n' + next_line
                    next_line = f.readline().strip()
                continue

        if content and slide_count_created < slide_count:
            slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
            title = slide.shapes.title
            title.text = header
            body_shape = slide.shapes.placeholders[slide_placeholder_index]
            tf = body_shape.text_frame
            tf.text = content


    prs.save(f'GeneratedPresentations/{ppt_name}.pptx')
    file_path = f"GeneratedPresentations/{ppt_name}.pptx"
    return f"{file_path}"


def generate_ppt(prompt, add_info, slides, theme, ppt_name):
    """
    Generate the PowerPoint presentation and save it.

    Args:
        prompt (str): The topic of the presentation.
        add_info (str): Additional information or instructions.
        slides (int): The number of slides.
        theme (int): The design theme number.
        ppt_name (str): The name of the generated PowerPoint file.

    Returns:
        str: The file path of the saved PowerPoint presentation.
    """
    theme = int(theme)

    # Validate and set theme
    if not theme or theme < 1 or theme > 7:
        print("Invalid theme number, default theme (1) will be applied.")
        theme = 1

    print("Generating the PowerPoint, this could take some time...\n")

    cache_dir = 'Cache'
    if not os.path.exists(cache_dir):
        os.makedirs(cache_dir)
    cache_file = os.path.join(cache_dir, f'{ppt_name}.txt')
    with open(cache_file, 'w', encoding='utf-8') as f:
        f.write(create_ppt_text(prompt, slides, add_info))

    # Create the PowerPoint presentation
    ppt_path = create_ppt(cache_file, theme, ppt_name,slides)
    return ppt_path

def save_presentation(response, ppt_name, theme, slides, additional_info):
    """
    Function to save the refined response into a PowerPoint presentation.

    Args:
        response (str): The refined response from the model.
        ppt_name (str): The desired name for the PowerPoint presentation.
        theme (int): The design theme number.
        slides (int): The number of slides.
        additional_info (str): Additional information or instructions.

    Returns:
        str: The file path of the saved PowerPoint presentation.
    """
    print("Saving the generated response as a PowerPoint presentation...\n")
    slides = slides + 2 # Because of Title and Table of Contents.


    # Generate the PowerPoint
    ppt_path = generate_ppt(response, additional_info, slides, theme, ppt_name)
    print(f"PowerPoint saved at: {ppt_path}")
    return ppt_path
